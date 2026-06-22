from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN_DARK = RGBColor(0x06, 0x4E, 0x3B)
GREEN_LIGHT = RGBColor(0x05, 0x96, 0x69)
CREAM = RGBColor(0xF8, 0xFA, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x11, 0x18, 0x27)
GRAY = RGBColor(0x6B, 0x72, 0x80)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        from lxml import etree
        solidFill = shape.fill._fill
        srgb = solidFill.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
        if srgb is not None:
            etree.SubElement(srgb, '{http://schemas.openxmlformats.org/drawingml/2006/main}alpha', val=str(int(alpha * 1000)))
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_slide(slide, left, top, width, height, items, font_size=18, color=DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return txBox

# ===== SLIDE 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

# Accent line
add_rect(slide, Inches(1), Inches(3.0), Inches(3), Inches(0.06), GREEN_LIGHT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.3),
    'KilimoLink Direct', 54, True, WHITE)
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(0.8),
    'Climate-Smart Food Transport Intelligence for Nairobi', 28, False, GREEN_LIGHT)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
    'Innovate4Cities 2026', 20, False, GRAY)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    'Gerald Kombo · github.com/geraldkombo/KilimoLink', 16, False, GRAY)

# ===== SLIDE 2: The Problem =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'The Climate Food Risk', 36, True, GREEN_DARK)

items = [
    'Nairobi\'s food comes from Kiambu, Machakos, Kajiado — all exposed to climate shocks',
    'Every truck kilometer adds emissions; empty return trips waste fuel',
    'Middlemen add cost, waste, and opacity to the food chain',
    'Kibera, Mathare, Mukuru face the sharpest price spikes and access shocks',
    'No real-time data exists to match supply zones to demand neighborhoods',
    'Food waste creates methane; long routes increase transport emissions',
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(8), Inches(5), items, 20, DARK)

# Stat box
add_rect(slide, Inches(9.5), Inches(1.6), Inches(3.2), Inches(2.2), GREEN_DARK)
add_text_box(slide, Inches(9.8), Inches(1.8), Inches(2.6), Inches(0.5),
    '40%', 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.8), Inches(2.4), Inches(2.6), Inches(1.2),
    'of food produced in Kenya is lost or wasted before reaching consumers', 14, False, CREAM, PP_ALIGN.CENTER)

# ===== SLIDE 3: Solution =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'The Solution: Climate-Smart Food Transport', 36, True, GREEN_DARK)

items = [
    'Connects farmers directly with Nairobi buyers — matching supply zones to demand neighborhoods',
    'Every transaction captures origin, destination, price, route, quantity',
    'Marketplace data becomes Nairobi\'s live food transport intelligence',
    'Reduces empty trips, shortens last-mile routes, cuts middlemen',
    'Estimates CO₂e emissions per delivery and per route',
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(8), Inches(4.5), items, 20, DARK)

# Flow boxes
flow_data = [
    ('Farmer\nLists Crop', GREEN_DARK),
    ('Buyer\nOrders', GREEN_LIGHT),
    ('Route\nOptimized', ORANGE),
    ('County\nDashboard', GREEN_DARK),
]
for i, (label, color) in enumerate(flow_data):
    x = Inches(9.0 + i * 1.05)
    y = Inches(2.0)
    add_rect(slide, x, y, Inches(0.9), Inches(1.2), color)
    add_text_box(slide, x, y + Inches(0.1), Inches(0.9), Inches(1.0),
        label, 11, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 4: County Dashboard =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'County Climate Intelligence Dashboard', 36, True, GREEN_DARK)

features = [
    ('Transport Flow Map', 'Which corridors are active, which are at risk'),
    ('Price Anomaly Heatmap', 'Real-time food inflation by neighborhood'),
    ('Supply Disruption Alerts', 'Weather + route + transaction signals'),
    ('Food Desert Index', 'Transport gap analysis for informal settlements'),
    ('CO₂e Impact Tiles', 'Per-route emissions savings estimates'),
]
for i, (title, desc) in enumerate(features):
    y = Inches(1.6 + i * 1.05)
    add_rect(slide, Inches(0.8), y, Inches(0.12), Inches(0.7), GREEN_LIGHT)
    add_text_box(slide, Inches(1.2), y, Inches(4), Inches(0.4),
        title, 20, True, GREEN_DARK)
    add_text_box(slide, Inches(1.2), y + Inches(0.4), Inches(6), Inches(0.4),
        desc, 16, False, GRAY)

# Demo callout
add_rect(slide, Inches(8.5), Inches(2.5), Inches(4), Inches(1.5), GREEN_DARK)
add_text_box(slide, Inches(8.8), Inches(2.7), Inches(3.4), Inches(1.2),
    'Live demo at\ngithub.com/geraldkombo/KilimoLink', 18, True, WHITE, PP_ALIGN.CENTER)

# ===== SLIDE 5: AI Engines =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'Three AI Engines', 36, True, GREEN_DARK)

ai_items = [
    ('Supply Disruption Predictor', 'Forecasts crop shortages by transport corridor and neighborhood using weather + transaction + route data'),
    ('Price Anomaly Detector', 'Flags unusual food inflation in informal settlements before families feel the shock'),
    ('Carbon Calculator', 'Estimates per-order emissions savings from shorter routes, fewer middlemen, avoided waste'),
]
colors = [GREEN_DARK, GREEN_LIGHT, ORANGE]
for i, (title, desc) in enumerate(ai_items):
    y = Inches(1.6 + i * 1.8)
    add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(1.5), colors[i])
    add_text_box(slide, Inches(1.2), y + Inches(0.15), Inches(10.5), Inches(0.5),
        f'0{i+1}  {title}', 24, True, WHITE)
    add_text_box(slide, Inches(1.2), y + Inches(0.75), Inches(10.5), Inches(0.6),
        desc, 16, False, CREAM)

# ===== SLIDE 6: Competition =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'Why KilimoLink Is Different', 36, True, GREEN_DARK)

# Table-like layout
headers = ['Alternative', 'What They Do', 'KilimoLink Difference']
rows = [
    ['Twiga Foods', 'B2B distribution', 'City climate + food transport data'],
    ['DigiFarm', 'Farmer services', 'Nairobi market + supply-demand zone matching'],
    ['M-Farm', 'Price information', 'Live transactions + route-based predictive alerts'],
    ['Generic marketplace', 'Connects buyers & sellers', 'Marketplace-as-data-engine for climate-smart logistics'],
]

y_start = 1.5
col_widths = [2.5, 3.5, 5.5]
for j, header in enumerate(headers):
    x = Inches(0.8 + sum(col_widths[:j]))
    add_text_box(slide, Inches(0.8 + sum(col_widths[:j])), Inches(y_start),
        Inches(col_widths[j]), Inches(0.5), header, 16, True, WHITE)
    add_rect(slide, Inches(0.8 + sum(col_widths[:j])), Inches(y_start),
        Inches(col_widths[j]), Inches(0.5), GREEN_DARK)

for i, row in enumerate(rows):
    y = y_start + 0.6 + i * 0.6
    bg = CREAM if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        add_rect(slide, Inches(0.8 + sum(col_widths[:j])), Inches(y),
            Inches(col_widths[j]), Inches(0.55), bg)
        add_text_box(slide, Inches(0.9 + sum(col_widths[:j])), Inches(y + 0.05),
            Inches(col_widths[j] - 0.2), Inches(0.45), cell, 14, j == 2, GREEN_DARK if j == 2 else DARK)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
    '"Competitors move food. KilimoLink helps Nairobi understand, predict, route, and climate-proof its food system."', 16, True, GREEN_LIGHT, PP_ALIGN.CENTER)

# ===== SLIDE 7: Informal Settlements =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'Informal Settlements Bonus Track', 36, True, GREEN_DARK)

items = [
    'Focus neighborhoods: Kibera, Mathare, Mukuru',
    'Real-time price monitoring for informal settlement food access',
    'Seller density and fresh produce availability tracking',
    'Shortage risk alerts before climate shocks hit vulnerable households',
    'Transport gap analysis linking supply zones to underserved neighborhoods',
]
add_bullet_slide(slide, Inches(0.8), Inches(1.6), Inches(8), Inches(4), items, 20, DARK)

add_rect(slide, Inches(9.5), Inches(1.6), Inches(3.2), Inches(2.5), GREEN_DARK)
add_text_box(slide, Inches(9.8), Inches(1.8), Inches(2.6), Inches(0.4),
    '60%', 44, True, WHITE, PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.8), Inches(2.4), Inches(2.6), Inches(1.2),
    'of Nairobi\'s population lives in informal settlements\n\n< 30%\nhave reliable access to fresh produce within 500m', 14, False, CREAM, PP_ALIGN.CENTER)

# ===== SLIDE 8: Roadmap =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), GREEN_DARK)

add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
    'Roadmap to August 31', 36, True, GREEN_DARK)

milestones = [
    ('Week 1-2', 'Fix deployment, stabilize demo, verify evidence citations', GREEN_DARK),
    ('Week 3-4', 'Ship county dashboard: transport flow map, heatmap, alert cards, impact tiles', GREEN_LIGHT),
    ('Week 5-6', 'Add AI disruption, anomaly, carbon engines + route risk maps', ORANGE),
    ('Week 7', 'Pilot supply-demand zone matching with farmers and buyers around Kibera, Mathare, Mukuru', GREEN_DARK),
    ('Week 8', 'Test, rehearse, and package Cambridge finale demo', GREEN_LIGHT),
]
for i, (week, desc, color) in enumerate(milestones):
    y = Inches(1.5 + i * 1.1)
    add_rect(slide, Inches(0.8), y, Inches(1.8), Inches(0.8), color)
    add_text_box(slide, Inches(0.8), y + Inches(0.15), Inches(1.8), Inches(0.5),
        week, 18, True, WHITE, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(3.0), y + Inches(0.1), Inches(9.5), Inches(0.6),
        desc, 18, False, DARK)

# ===== SLIDE 9: Ask =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, GREEN_DARK)

add_rect(slide, Inches(1), Inches(2.8), Inches(3), Inches(0.06), GREEN_LIGHT)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    'The Ask', 48, True, WHITE)
add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
    'Pilot with Nairobi City County to map, route, and climate-proof the city\'s food system.', 24, False, GREEN_LIGHT)

items = [
    'Access to county GIS and market datasets',
    'Support to onboard informal settlement buyers and peri-urban farmers',
    'Help scale the food transport intelligence dashboard city-wide',
]
add_bullet_slide(slide, Inches(1), Inches(4.0), Inches(10), Inches(2.5), items, 20, CREAM)

# Closing line
add_rect(slide, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0), RGBColor(0x04, 0x3B, 0x2D))
add_text_box(slide, Inches(1), Inches(6.6), Inches(11), Inches(0.8),
    'KilimoLink Direct helps Nairobi see where food is, which routes carry it, where climate risk will hit next — before families feel the shock.', 18, True, WHITE, PP_ALIGN.CENTER)

output_path = 'C:\\Users\\Rosemary\\Desktop\\AgriBizPlatform_FINAL\\I4C26_PITCH_DECK.pptx'
prs.save(output_path)
print(f'Deck saved to {output_path}')
